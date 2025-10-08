"""PowerPoint Generator - Simple version for maximum compatibility."""

from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .llm_summary_generator import LLMSummaryGenerator


class PowerPointGenerator:
    """Creates PowerPoint presentations with sprint summaries - simplified for compatibility."""

    # Color scheme
    COLORS = {
        "good": RGBColor(76, 175, 80),      # Green
        "fair": RGBColor(255, 167, 38),     # Orange
        "poor": RGBColor(239, 83, 80),      # Red
        "gray": RGBColor(117, 117, 117),    # Gray
        "dark_gray": RGBColor(54, 54, 54),  # Dark gray
        "light_gray": RGBColor(245, 245, 245),  # Light gray background
        "border": RGBColor(221, 221, 221),  # Border gray
        "text": RGBColor(66, 66, 66),       # Text gray
        "blocker_red": RGBColor(211, 47, 47),  # Blocker red
        "success_green": RGBColor(56, 142, 60),  # Success green
        "blue": RGBColor(33, 150, 243),     # Blue accent
    }

    def __init__(self, provider: str, api_key: str, model: str):
        """Initialize with LLM provider configuration."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.llm_generator = LLMSummaryGenerator(provider, api_key, model)

    def generate_presentation(
        self,
        all_summaries: List[Dict[str, Any]],
        all_sprint_data: List[Dict[str, Any]],
        all_metrics: List[Dict[str, Any]],
        output_dir: str = "./output",
    ):
        """Generate PowerPoint presentation with LLM-powered content."""
        # Add title slide
        self._create_title_slide(all_summaries)

        # Add slide for each team
        for i, (summary, sprint_data, metrics) in enumerate(zip(all_summaries, all_sprint_data, all_metrics)):
            print(f"   🤖 Generating AI slide content for {summary['teamInfo']['label']}...")
            self._create_team_slide(summary, sprint_data, metrics)

        # Save presentation
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        file_path = output_path / "sprint-summary-presentation.pptx"

        self.prs.save(str(file_path))
        print(f"   ✅ PowerPoint saved to: {file_path}")

    def _create_title_slide(self, all_summaries: List[Dict[str, Any]]):
        """Create title slide - simplified version."""
        title_layout = self.prs.slide_layouts[0]  # Use title slide layout
        slide = self.prs.slides.add_slide(title_layout)

        # Set title
        title = slide.shapes.title
        title.text = "📊 Sprint Summary Report"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.COLORS["blue"]

        # Set subtitle with projects and teams
        projects = sorted(set(s["projectInfo"]["key"] for s in all_summaries))
        teams = sorted(set(s["teamInfo"]["label"] for s in all_summaries))

        subtitle = slide.placeholders[1]
        subtitle.text = f"📁 Projects: {', '.join(projects)}\n👥 Teams: {', '.join(teams)}\n\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self.COLORS["dark_gray"]

    def _create_team_slide(
        self,
        summary: Dict[str, Any],
        sprint_data: Dict[str, Any],
        metrics: Dict[str, Any],
    ):
        """Create a slide for a specific team with 2x2 layout - simplified."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        team_label = summary["teamInfo"]["label"]
        project_key = summary["projectInfo"]["key"]
        health = summary["sprintHealthMetrics"]["overallHealth"]

        # Title at top
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = f"Team: {team_label} ({project_key}) - Health: {health}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self._get_health_color(health)

        # Get LLM-generated content
        llm_content = self.llm_generator.generate_slide_content(summary)

        # 2x2 layout positions - use simple textboxes only
        positions = [
            (Inches(0.5), Inches(1.2)),  # Top left
            (Inches(5.2), Inches(1.2)),  # Top right
            (Inches(0.5), Inches(4.2)),  # Bottom left
            (Inches(5.2), Inches(4.2)),  # Bottom right
        ]

        box_width = Inches(4.3)
        box_height = Inches(2.8)

        # Top Left: Health Summary
        self._add_simple_box(
            slide,
            positions[0][0],
            positions[0][1],
            box_width,
            box_height,
            llm_content["healthSummary"]["title"],
            llm_content["healthSummary"]["bullets"],
            self._get_health_color(health),
        )

        # Top Right: Accomplishments
        self._add_simple_box(
            slide,
            positions[1][0],
            positions[1][1],
            box_width,
            box_height,
            llm_content["accomplishments"]["title"],
            llm_content["accomplishments"]["bullets"],
            self.COLORS["success_green"],
        )

        # Bottom Left: Blockers
        blocker_color = (
            self.COLORS["blocker_red"]
            if summary.get("currentBlockers") and len(summary["currentBlockers"]) > 0
            else self.COLORS["success_green"]
        )
        self._add_simple_box(
            slide,
            positions[2][0],
            positions[2][1],
            box_width,
            box_height,
            llm_content["blockers"]["title"],
            llm_content["blockers"]["bullets"],
            blocker_color,
        )

        # Bottom Right: Recommendations
        self._add_simple_box(
            slide,
            positions[3][0],
            positions[3][1],
            box_width,
            box_height,
            llm_content["recommendations"]["title"],
            llm_content["recommendations"]["bullets"],
            self.COLORS["blue"],
        )

    def _add_simple_box(
        self,
        slide,
        left,
        top,
        width,
        height,
        title,
        bullets,
        title_color,
    ):
        """Add a simple text box - no shapes, just text for compatibility."""
        # Create textbox for entire area
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Pt(10)
        text_frame.margin_left = Pt(10)
        text_frame.margin_right = Pt(10)

        # Add title
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_after = Pt(8)

        # Add bullets
        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.COLORS["text"]
            p.space_before = Pt(4)

    def _get_health_color(self, health: str) -> RGBColor:
        """Get color for health status."""
        health_lower = health.lower()
        if health_lower == "good":
            return self.COLORS["good"]
        elif health_lower == "fair":
            return self.COLORS["fair"]
        elif health_lower == "poor":
            return self.COLORS["poor"]
        else:
            return self.COLORS["gray"]
