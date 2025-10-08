"""PowerPoint Generator - Creates presentation slides with sprint summaries."""

from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .llm_summary_generator import LLMSummaryGenerator


class PowerPointGenerator:
    """Creates PowerPoint presentations with sprint summaries in 2x2 format."""

    # Color scheme
    COLORS = {
        "good": RGBColor(76, 175, 80),      # Green
        "fair": RGBColor(255, 167, 38),     # Orange
        "poor": RGBColor(239, 83, 80),      # Red
        "gray": RGBColor(117, 117, 117),    # Gray
        "dark_gray": RGBColor(54, 54, 54),  # Dark gray
        "light_gray": RGBColor(245, 245, 245),  # Light gray background
        "border": RGBColor(221, 221, 221),  # Border gray
        "text": RGBColor(66, 66, 66),       # Text gray
        "blocker_red": RGBColor(211, 47, 47),  # Blocker red
        "success_green": RGBColor(56, 142, 60),  # Success green
        "blue": RGBColor(33, 150, 243),     # Blue accent
    }

    def __init__(self, provider: str, api_key: str, model: str):
        """Initialize with LLM provider configuration."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.llm_generator = LLMSummaryGenerator(provider, api_key, model)

    def generate_presentation(
        self,
        all_summaries: List[Dict[str, Any]],
        all_sprint_data: List[Dict[str, Any]],
        all_metrics: List[Dict[str, Any]],
        output_dir: str = "./output",
    ):
        """Generate PowerPoint presentation with LLM-powered content."""
        # Add title slide
        self._create_title_slide(all_summaries)

        # Add slide for each team
        for i, (summary, sprint_data, metrics) in enumerate(zip(all_summaries, all_sprint_data, all_metrics)):
            print(f"   🤖 Generating AI slide content for {summary['teamInfo']['label']}...")
            self._create_team_slide(summary, sprint_data, metrics)

        # Save presentation
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        file_path = output_path / "sprint-summary-presentation.pptx"

        self.prs.save(str(file_path))
        print(f"   ✅ PowerPoint saved to: {file_path}")

    def _create_title_slide(self, all_summaries: List[Dict[str, Any]]):
        """Create title slide with enhanced graphics."""
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)

        # Add colorful header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(1.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.COLORS["blue"]

        # Main title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "📊 Sprint Summary Report"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Projects and teams info
        projects = sorted(set(s["projectInfo"]["key"] for s in all_summaries))
        teams = sorted(set(s["teamInfo"]["label"] for s in all_summaries))

        # Info box
        info_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(2.5), Inches(6), Inches(2.5)
        )
        info_box.fill.solid()
        info_box.fill.fore_color.rgb = self.COLORS["light_gray"]
        info_box.line.color.rgb = self.COLORS["border"]
        info_box.line.width = Pt(1)

        # Info text
        info_text = slide.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(5), Inches(1.5))
        info_frame = info_text.text_frame
        info_frame.word_wrap = True

        # Projects
        p1 = info_frame.paragraphs[0]
        p1.text = f"📁 Projects: {', '.join(projects)}"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = self.COLORS["dark_gray"]

        # Teams
        p2 = info_frame.add_paragraph()
        p2.text = f"👥 Teams: {', '.join(teams)}"
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = self.COLORS["dark_gray"]
        p2.space_before = Pt(12)

        # Generated date
        date_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(4), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = self.COLORS["gray"]
        date_para.alignment = PP_ALIGN.CENTER

    def _create_team_slide(
        self,
        summary: Dict[str, Any],
        sprint_data: Dict[str, Any],
        metrics: Dict[str, Any],
    ):
        """Create a slide for a specific team with 2x2 layout."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        team_label = summary["teamInfo"]["label"]
        project_key = summary["projectInfo"]["key"]
        health = summary["sprintHealthMetrics"]["overallHealth"]

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = f"Team: {team_label} ({project_key})"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS["dark_gray"]

        # Health indicator next to title
        health_color = self._get_health_color(health)
        health_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(8.5), Inches(0.25), Inches(0.5), Inches(0.5)
        )
        health_circle.fill.solid()
        health_circle.fill.fore_color.rgb = health_color

        # Get LLM-generated content using the summary which already has all the data
        llm_content = self.llm_generator.generate_slide_content(summary)

        # 2x2 layout positions
        positions = [
            (Inches(0.5), Inches(1.2)),  # Top left
            (Inches(5.2), Inches(1.2)),  # Top right
            (Inches(0.5), Inches(4.2)),  # Bottom left
            (Inches(5.2), Inches(4.2)),  # Bottom right
        ]

        box_width = Inches(4.3)
        box_height = Inches(2.8)

        # Top Left: Health Summary
        self._add_enhanced_box(
            slide,
            positions[0][0],
            positions[0][1],
            box_width,
            box_height,
            llm_content["healthSummary"]["title"],
            llm_content["healthSummary"]["bullets"],
            self.COLORS["light_gray"],
            self.COLORS["border"],
            health_indicator=health_color,
        )

        # Top Right: Accomplishments
        self._add_enhanced_box(
            slide,
            positions[1][0],
            positions[1][1],
            box_width,
            box_height,
            llm_content["accomplishments"]["title"],
            llm_content["accomplishments"]["bullets"],
            self.COLORS["light_gray"],
            self.COLORS["border"],
        )

        # Bottom Left: Blockers
        blocker_color = (
            self.COLORS["blocker_red"]
            if summary.get("currentBlockers") and len(summary["currentBlockers"]) > 0
            else self.COLORS["success_green"]
        )
        self._add_enhanced_box(
            slide,
            positions[2][0],
            positions[2][1],
            box_width,
            box_height,
            llm_content["blockers"]["title"],
            llm_content["blockers"]["bullets"],
            self.COLORS["light_gray"],
            self.COLORS["border"],
            text_color=blocker_color,
        )

        # Bottom Right: Recommendations
        self._add_enhanced_box(
            slide,
            positions[3][0],
            positions[3][1],
            box_width,
            box_height,
            llm_content["recommendations"]["title"],
            llm_content["recommendations"]["bullets"],
            self.COLORS["light_gray"],
            self.COLORS["border"],
        )

    def _add_enhanced_box(
        self,
        slide,
        left,
        top,
        width,
        height,
        title,
        bullets,
        fill_color,
        border_color,
        health_indicator=None,
        text_color=None,
    ):
        """Add an enhanced info box with rounded corners and optional health indicator."""
        # Background rectangle with rounded corners
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
        box.line.color.rgb = border_color
        box.line.width = Pt(1)

        # Title
        title_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(14)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS["dark_gray"]

        # Health indicator circle (if provided)
        if health_indicator:
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                left + width - Inches(0.7),
                top + Inches(0.15),
                Inches(0.5),
                Inches(0.5),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = health_indicator

        # Bullets
        bullets_box = slide.shapes.add_textbox(
            left + Inches(0.1),
            top + Inches(0.6),
            width - Inches(0.2),
            height - Inches(0.7),
        )
        bullets_frame = bullets_box.text_frame
        bullets_frame.word_wrap = True
        bullets_frame.vertical_anchor = MSO_ANCHOR.TOP

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = bullets_frame.paragraphs[0]
            else:
                p = bullets_frame.add_paragraph()

            p.text = f"• {bullet}"
            p.font.size = Pt(11)
            p.font.color.rgb = text_color if text_color else self.COLORS["text"]
            p.space_before = Pt(4) if i > 0 else Pt(0)

    def _get_health_color(self, health: str) -> RGBColor:
        """Get color for health status."""
        health_lower = health.lower()
        if health_lower == "good":
            return self.COLORS["good"]
        elif health_lower == "fair":
            return self.COLORS["fair"]
        elif health_lower == "poor":
            return self.COLORS["poor"]
        else:
            return self.COLORS["gray"]
