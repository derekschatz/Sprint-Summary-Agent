#!/usr/bin/env python3
"""Create ultra-compatible PowerPoint using pptx-python with explicit compatibility settings."""

import sys
sys.path.insert(0, '/Users/derekschatz/Documents/GitHub/Sprint-Summary-Agent')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json

# Load the summary data
with open('output/sprint-summary-combined.json', 'r') as f:
    data = json.load(f)

# Create presentation with explicit defaults
prs = Presentation()

# Use Title Slide layout for first slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

# Set title
title = slide.shapes.title
title.text = "Sprint Summary Report"

# Set subtitle
if len(slide.placeholders) > 1:
    subtitle = slide.placeholders[1]
    projects = ", ".join([p["key"] for p in data["projects"]])
    teams = ", ".join(data["teams"])
    subtitle.text = f"Projects: {projects}\nTeams: {teams}\n\nGenerated: {data['generatedAt'][:10]}"

# Add one team slide as example
if data.get("teamSummaries"):
    team = data["teamSummaries"][0]

    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Simple title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = f"Team: {team['team']} - {team['health']}"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True

    # Simple content box
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = f"Completion Rate: {team['completionRate']}"
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = f"Velocity: {team['velocity']} issues completed"
    p.font.size = Pt(18)

# Save with minimal settings
output_file = "output/sprint-summary-ultra-compatible.pptx"
prs.save(output_file)
print(f"Created ultra-compatible PowerPoint: {output_file}")
print(f"Please try opening this file and let me know if it works.")
