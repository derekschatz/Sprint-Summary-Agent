#!/usr/bin/env python3
"""Fix PowerPoint compatibility by creating a new file with better compatibility."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def fix_powerpoint_compatibility():
    """Create a new PowerPoint with maximum compatibility."""
    # Load the existing presentation
    prs_old = Presentation("output/sprint-summary-presentation.pptx")

    # Create a new presentation with default template (better compatibility)
    prs = Presentation()

    # Copy slides one by one with simpler structure
    for slide_num, old_slide in enumerate(prs_old.slides):
        # Use blank layout
        blank_layout = prs.slide_layouts[6]
        new_slide = prs.slides.add_slide(blank_layout)

        # Copy all shapes
        for shape in old_slide.shapes:
            # Get shape properties
            el = shape.element
            new_el = el
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # Save with compatibility mode
    output_path = "output/sprint-summary-presentation-compatible.pptx"
    prs.save(output_path)
    print(f"Created compatible version: {output_path}")

if __name__ == "__main__":
    fix_powerpoint_compatibility()
